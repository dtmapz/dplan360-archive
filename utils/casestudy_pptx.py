"""Case Study 16:9 슬라이드 PPTX 렌더러 (python-pptx).

utils/casestudy_render.py 의 HTML 레이아웃을 동일한 좌표계(1280×720 px)로 재현.
1 px = 9525 EMU.
"""
import io
import re
import requests
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


PX = 9525  # EMU per px

# Colors
NAVY = RGBColor(0x0F, 0x1E, 0x3D)
BLUE = RGBColor(0x4C, 0x7D, 0xFF)
BLUE_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
GRAY_BG = RGBColor(0xF5, 0xF7, 0xFB)
GRAY_BORDER = RGBColor(0xE5, 0xE9, 0xF0)
GRAY_DIVIDER = RGBColor(0xD1, 0xD5, 0xDB)
TEXT_PRIMARY = NAVY
TEXT_BODY = RGBColor(0x37, 0x41, 0x51)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
HEADER_SUB = RGBColor(0xA5, 0xB4, 0xCE)
SCOPE_INTERNAL = RGBColor(0x35, 0x47, 0x6B)
SCOPE_INTERNAL_FG = RGBColor(0xE5, 0xE9, 0xF0)
SCOPE_EXTERNAL = RGBColor(0x16, 0xA3, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREATIVE_BG_1 = RGBColor(0xD4, 0xC5, 0xB0)

FONT_NAME = "Pretendard"  # macOS 없으면 시스템 fallback


# ---------------------------------------------------------------------
# 저수준 헬퍼
# ---------------------------------------------------------------------

def _emu(x_px: float) -> Emu:
    return Emu(int(x_px * PX))


def _rect(shapes, x, y, w, h, fill=None, line=None):
    shp = shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(x), _emu(y), _emu(w), _emu(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Emu(int(0.5 * PX))
    # remove default text placeholder styling
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return shp


def _textbox(shapes, x, y, w, h, text, *, size=12, bold=False, color=TEXT_BODY,
             align="left", anchor="top", spacing=None, font=FONT_NAME):
    tb = shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    _add_run(p, text, size=size, bold=bold, color=color, spacing=spacing, font=font)
    return tb


def _add_run(paragraph, text, *, size=12, bold=False, color=TEXT_BODY,
             spacing=None, font=FONT_NAME):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if spacing is not None:
        # letter-spacing in Pt (approx via spc XML attr, 100ths of a point)
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return run


def _pill(shapes, x, y, text, *, size=10.5, bold=True, bg=BLUE_LIGHT, fg=BLUE,
          padding_x=8, padding_y=3, spacing=0.3, radius_px=3):
    """Auto-width pill. width = text_len * size * 0.6 + 2*padding_x (approx)."""
    text = str(text)
    est_w = len(text) * size * 0.62 + 2 * padding_x
    h = size * 1.3 + 2 * padding_y
    shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _emu(x), _emu(y), _emu(est_w), _emu(h))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    # rounded rectangle small radius: adjust adjustment value (0.0~0.5)
    try:
        shp.adjustments[0] = radius_px / (h / 2)
    except Exception:
        pass
    tf = shp.text_frame
    tf.margin_left = _emu(padding_x)
    tf.margin_right = _emu(padding_x)
    tf.margin_top = _emu(padding_y - 1)
    tf.margin_bottom = _emu(padding_y - 1)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, text, size=size, bold=bold, color=fg, spacing=spacing)
    return shp, est_w


def _line_h(shapes, x, y, w, color):
    shp = shapes.add_connector(1, _emu(x), _emu(y), _emu(x + w), _emu(y))
    shp.line.color.rgb = color
    shp.line.width = Emu(int(1 * PX))
    return shp


def _line_v(shapes, x, y, h, color):
    shp = shapes.add_connector(1, _emu(x), _emu(y), _emu(x), _emu(y + h))
    shp.line.color.rgb = color
    shp.line.width = Emu(int(1 * PX))
    return shp


# ---------------------------------------------------------------------
# 이미지 로딩
# ---------------------------------------------------------------------

def _fetch_image_bytes(url: str) -> bytes | None:
    if not url:
        return None
    try:
        r = requests.get(url, timeout=8)
        if r.status_code == 200 and r.headers.get("Content-Type", "").startswith("image/"):
            return r.content
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------
# 파싱 헬퍼
# ---------------------------------------------------------------------

_BRACKET_RE = re.compile(r"(\[[^\]]+\])")


def _title_segments(title: str) -> list[tuple[str, bool]]:
    """[대괄호]로 감싼 부분을 accent=True 로 분리 + \n 을 line-break 로."""
    segments: list[tuple[str, bool]] = []
    for line_idx, line in enumerate(title.split("\n")):
        if line_idx > 0:
            segments.append(("\n", False))
        for part in _BRACKET_RE.split(line):
            if not part:
                continue
            if part.startswith("[") and part.endswith("]"):
                segments.append((part[1:-1], True))
            else:
                segments.append((part, False))
    return segments


def _campaign_type_pills_text(types: list[str]) -> list[str]:
    types = [t.strip() for t in (types or []) if t.strip()]
    if not types:
        return []
    if len(types) >= 3:
        return ["FULL-FUNNEL"]
    return [t.upper() for t in types]


def _period_str(s, e) -> str:
    s, e = str(s or "").strip(), str(e or "").strip()
    if s and e and s != e:
        return f"{s} ~ {e}"
    return s or e or "-"


def _target_str(g, a) -> str:
    g, a = str(g or "").strip(), str(a or "").strip()
    if g and a:
        return f"{g} · {a}"
    return g or a or "-"


# ---------------------------------------------------------------------
# 섹션 렌더러
# ---------------------------------------------------------------------

def _render_header(shapes, cs: dict):
    _rect(shapes, 0, 0, 1280, 72, fill=NAVY)

    # CASE STUDY badge
    badge_shp, badge_w = _pill(
        shapes, 40, 23, "CASE STUDY",
        size=11, bg=BLUE, fg=WHITE, padding_x=12, padding_y=6, spacing=1.2, radius_px=3,
    )
    # brand info
    brand_x = 40 + badge_w + 20
    _textbox(shapes, brand_x, 23, 400, 30,
             cs.get("brand", ""), size=20, bold=True, color=WHITE, anchor="middle")
    right_meta = " · ".join(x for x in [cs.get("advertiser"), cs.get("industry")] if x)
    _textbox(shapes, brand_x + 200, 27, 400, 24,
             right_meta, size=13, color=HEADER_SUB, anchor="middle")

    # Right: scope badge + logo
    scope = (cs.get("share_scope") or "Internal").strip()
    is_ext = scope.lower() == "external"
    scope_bg = SCOPE_EXTERNAL if is_ext else SCOPE_INTERNAL
    scope_fg = WHITE if is_ext else SCOPE_INTERNAL_FG
    scope_label = scope.upper()

    logo_w = 90
    logo_x = 1280 - 40 - logo_w
    _textbox(shapes, logo_x, 27, logo_w, 24,
             "D-PLAN360", size=12, bold=True, color=HEADER_SUB,
             align="right", anchor="middle", spacing=2.0)

    scope_pill_est = len(scope_label) * 10.5 * 0.62 + 20 + 10
    scope_x = logo_x - 16 - scope_pill_est
    _pill(shapes, scope_x, 27, scope_label,
          size=10.5, bg=scope_bg, fg=scope_fg,
          padding_x=10, padding_y=5, spacing=1.2, radius_px=3)


def _render_meta_strip(shapes, cs: dict):
    _rect(shapes, 0, 72, 1280, 40, fill=GRAY_BG)
    _line_h(shapes, 0, 112, 1280, GRAY_BORDER)

    items = [
        ("MEDIA", cs.get("media", "-")),
        ("PERIOD", _period_str(cs.get("period_start"), cs.get("period_end"))),
        ("TARGET", _target_str(cs.get("target_gender"), cs.get("target_age"))),
    ]

    x = 40
    y_base = 72
    row_h = 40

    for i, (key, val) in enumerate(items):
        key_w = len(key) * 10.5 * 0.65 + 4
        _textbox(shapes, x, y_base, key_w, row_h,
                 key, size=10.5, bold=True, color=TEXT_MUTED,
                 anchor="middle", spacing=1.0)
        x += key_w + 8
        val_w = len(str(val)) * 12 * 0.65 + 10
        _textbox(shapes, x, y_base, val_w, row_h,
                 str(val), size=12, bold=True, color=TEXT_PRIMARY, anchor="middle")
        x += val_w + 24
        if i < len(items):
            _line_v(shapes, x - 12, y_base + 13, 14, GRAY_DIVIDER)

    # TYPE
    key = "TYPE"
    key_w = len(key) * 10.5 * 0.65 + 4
    _textbox(shapes, x, y_base, key_w, row_h,
             key, size=10.5, bold=True, color=TEXT_MUTED,
             anchor="middle", spacing=1.0)
    x += key_w + 8
    for label in _campaign_type_pills_text(cs.get("campaign_types", [])):
        pill_y = y_base + (row_h - (10.5 * 1.3 + 6)) / 2
        _, pw = _pill(shapes, x, pill_y, label,
                      size=10.5, bg=BLUE_LIGHT, fg=BLUE,
                      padding_x=8, padding_y=3, spacing=0.3, radius_px=3)
        x += pw + 4


def _render_creative(shapes, cs: dict):
    # 560×315 영역 at x=40, y=132. 이미지 1/2/4개를 동일 영역 안에서 분할 배치.
    x, y, w, h = 40, 132, 560, 315
    gap = 2

    urls = cs.get("creative_image_urls") or (
        [cs["creative_image_url"]] if cs.get("creative_image_url") else []
    )
    urls = [u for u in urls if u]

    if not urls:
        _rect(shapes, x, y, w, h, fill=CREATIVE_BG_1)
        _textbox(shapes, x, y, w, h,
                 "CAMPAIGN CREATIVE · 16:9", size=13, color=WHITE,
                 align="center", anchor="middle", spacing=1.5)
        return

    layout_n = 1 if len(urls) == 1 else (2 if len(urls) == 2 else 4)
    urls = urls[:layout_n]

    if layout_n == 1:
        boxes = [(x, y, w, h)]
    elif layout_n == 2:
        tw = (w - gap) / 2
        boxes = [(x, y, tw, h), (x + tw + gap, y, tw, h)]
    else:
        tw = (w - gap) / 2
        th = (h - gap) / 2
        boxes = [
            (x, y, tw, th), (x + tw + gap, y, tw, th),
            (x, y + th + gap, tw, th), (x + tw + gap, y + th + gap, tw, th),
        ]

    for (bx, by, bw, bh), url in zip(boxes, urls):
        img_bytes = _fetch_image_bytes(url)
        if img_bytes:
            try:
                shapes.add_picture(io.BytesIO(img_bytes), _emu(bx), _emu(by),
                                   width=_emu(bw), height=_emu(bh))
                continue
            except Exception:
                pass
        _rect(shapes, bx, by, bw, bh, fill=CREATIVE_BG_1)


def _render_caption(shapes, ai: dict):
    _textbox(shapes, 42, 455, 560, 30,
             ai.get("caption", ""), size=13, color=TEXT_MUTED, anchor="top")


def _render_title_and_results(shapes, cs: dict, ai: dict):
    x, y = 632, 132
    w = 1280 - 40 - x  # 608

    # eyebrow
    _textbox(shapes, x, y, w, 18,
             ai.get("eyebrow", ""), size=11.5, bold=True, color=BLUE, spacing=1.8)

    # title (multi-line, accent segments in blue)
    title_tb = shapes.add_textbox(_emu(x), _emu(y + 26), _emu(w), _emu(100))
    tf = title_tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    para.line_spacing = 1.28
    first_run = True
    for text, accent in _title_segments(ai.get("title", "")):
        if text == "\n":
            para = tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            para.line_spacing = 1.28
            first_run = True
            continue
        color = BLUE if accent else NAVY
        run = para.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(23)
        run.font.bold = True
        run.font.color.rgb = color
        first_run = False

    # Results box
    rx, ry, rw, rh = x, 330, w, 130
    _rect(shapes, rx, ry, rw, rh, fill=GRAY_BG)
    _rect(shapes, rx, ry, 3, rh, fill=BLUE)  # left border

    _textbox(shapes, rx + 18, ry + 12, rw - 36, 16,
             "RESULTS", size=11, bold=True, color=BLUE, spacing=1.5)

    # metrics
    results = [r for r in (cs.get("results") or []) if r.get("kpi_name") or r.get("value")][:4]
    n = max(len(results), 1)
    gap = 12
    inner_x = rx + 18
    inner_w = rw - 36
    cell_w = (inner_w - gap * (n - 1)) / n
    cell_y = ry + 42
    for i, r in enumerate(results):
        cx = inner_x + i * (cell_w + gap)
        _textbox(shapes, cx, cell_y, cell_w, 16,
                 r.get("kpi_name", ""), size=11.5, color=TEXT_MUTED)
        is_hero = (i == 0)
        _textbox(shapes, cx, cell_y + 20, cell_w, 36,
                 r.get("value", ""),
                 size=(26 if is_hero else 24), bold=True,
                 color=(BLUE if is_hero else NAVY))


def _render_bottom_sections(shapes, ai: dict):
    y = 500
    total_w = 1280 - 80
    gap = 28
    col_w = (total_w - gap * 2) / 3

    sections = [
        ("Challenge", "캠페인 목표", ai.get("challenge_bullets", []) or []),
        ("Approach", "캠페인 전략", ai.get("approach_bullets", []) or []),
        ("Insight", "인사이트 · 테스트", ai.get("insight_bullets", []) or []),
    ]

    for i, (en, kr, bullets) in enumerate(sections):
        x = 40 + i * (col_w + gap)
        # title left (en)
        _textbox(shapes, x, y, col_w * 0.55, 20,
                 en, size=15, bold=True, color=NAVY)
        # title right (kr)
        _textbox(shapes, x + col_w * 0.55, y + 3, col_w * 0.45, 18,
                 kr, size=12, color=TEXT_MUTED, align="right")
        # underline
        _line_h(shapes, x, y + 26, col_w, NAVY)

        # bullets
        by = y + 36
        for b in bullets[:5]:
            # dot
            dot = shapes.add_shape(MSO_SHAPE.OVAL, _emu(x), _emu(by + 6),
                                   _emu(5), _emu(5))
            dot.shadow.inherit = False
            dot.fill.solid(); dot.fill.fore_color.rgb = BLUE
            dot.line.fill.background()
            # text
            tb = _textbox(shapes, x + 12, by, col_w - 12, 40,
                          str(b), size=14, color=TEXT_BODY)
            tb.text_frame.paragraphs[0].line_spacing = 1.5
            # advance y by measured lines (approx)
            approx_lines = max(1, int(len(str(b)) * 14 * 0.62 / (col_w - 12)) + 1)
            by += approx_lines * 21 + 6


def _render_footer(shapes):
    _rect(shapes, 0, 680, 1280, 40, fill=GRAY_BG)
    _line_h(shapes, 0, 680, 1280, GRAY_BORDER)


# ---------------------------------------------------------------------
# 진입점
# ---------------------------------------------------------------------

def build_slide_pptx(cs: dict, ai: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = _emu(1280)
    prs.slide_height = _emu(720)

    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    shapes = slide.shapes

    # white background
    _rect(shapes, 0, 0, 1280, 720, fill=WHITE)

    _render_header(shapes, cs)
    _render_meta_strip(shapes, cs)
    _render_creative(shapes, cs)
    _render_caption(shapes, ai)
    _render_title_and_results(shapes, cs, ai)
    _render_bottom_sections(shapes, ai)
    _render_footer(shapes)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
